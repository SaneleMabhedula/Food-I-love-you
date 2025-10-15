from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import requests
from io import BytesIO

# Create a presentation
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "South African Traditional Attires"
subtitle.text = "A glimpse of vibrant cultural fashion"

# Content Slide
slide_layout = prs.slide_layouts[5]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Add title
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "South African Attire"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)

# Add description
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(2))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = ("South African traditional attire is colorful and diverse, "
          "reflecting the rich heritage of different cultures such as "
          "Zulu, Xhosa, Sotho, and Ndebele. These outfits are often worn "
          "during weddings, ceremonies, and celebrations, showcasing beadwork, "
          "patterns, and unique designs.")

# Insert an image from the web (royalty-free)
img_url = "https://upload.wikimedia.org/wikipedia/commons/5/59/Zulu_women_in_traditional_dress.jpg"
response = requests.get(img_url)
img_stream = BytesIO(response.content)
slide.shapes.add_picture(img_stream, Inches(1.5), Inches(3), Inches(6), Inches(3.5))

# Save the presentation
prs.save("South_Africa_Attire_Presentation.pptx")
print("Presentation created successfully!")
